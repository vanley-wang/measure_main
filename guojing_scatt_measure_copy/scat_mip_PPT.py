import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN

# === 路径配置 ===
path_0701 = 'Data/FXN_2023/FXN_20230701/scatt_mip'
path_0703 = 'Data/FXN_2023/FXN_20230703/scatt_mip'
output_pptx = 'Data/FXN_2023/散射系数投影图.pptx'

# === 图像收集函数 ===
def collect_images(path):
    img_dict = {}
    for fname in os.listdir(path):
        if not fname.endswith('.png'):
            continue
        parts = fname.split('_')
        if len(parts) < 4:
            continue
        sample = parts[0]
        time = parts[1]
        direction = parts[-1].replace('.png', '')  # Z / Y
        key = f"{sample}_{direction}_{time}"
        img_dict[key] = os.path.join(path, fname)
    return img_dict

img_dict = {}
img_dict.update(collect_images(path_0701))
img_dict.update(collect_images(path_0703))

# === 样本列表 ===
samples = sorted(set(key.split('_')[0] for key in img_dict if '_Z_0701' in key))

# === PPT 初始化 ===
prs = Presentation()
blank_slide = prs.slide_layouts[6]

def insert_sample(slide, sample, direction='left'):
    # 图像尺寸
    img_width_z, img_height_z = Cm(6.0), Cm(6.0)
    img_width_y, img_height_y = Cm(6.0), Cm(3.75)
    gap = Cm(0.25)
    shift_y = Cm(3.5)  # 图片整体下移

    # 坐标设置
    x0 = Cm(0.25) if direction == 'left' else Cm(12.75)
    title_top = Cm(0.5)
    date_top = Cm(3.2)
    z_top = Cm(5.0)
    y_top = Cm(13.0)

    # 样本标题
    title_box = slide.shapes.add_textbox(x0, title_top, width=img_width_z * 2 + gap, height=Cm(0.8))
    tf = title_box.text_frame
    tf.text = sample
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 日期行（0701 和 0703），各居中于 Z 图与 Y 图的两列整体
    for i, date in enumerate(['0701', '0703']):
        x_date = x0 + i * (img_width_z + gap)
        date_box = slide.shapes.add_textbox(x_date, date_top, width=img_width_z, height=Cm(0.5))
        tf = date_box.text_frame
        tf.text = date
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 图像路径
    img_z_0701 = img_dict.get(f"{sample}_Z_0701")
    img_z_0703 = img_dict.get(f"{sample}_Z_0703")
    img_y_0701 = img_dict.get(f"{sample}_Y_0701")
    img_y_0703 = img_dict.get(f"{sample}_Y_0703")

    if not all([img_z_0701, img_z_0703, img_y_0701, img_y_0703]):
        print(f"⚠️ 缺图跳过样本: {sample}")
        return

    # 插入 Z 图
    slide.shapes.add_picture(img_z_0701, x0, z_top, width=img_width_z, height=img_height_z)
    slide.shapes.add_picture(img_z_0703, x0 + img_width_z + gap, z_top, width=img_width_z, height=img_height_z)

    # 插入 Y 图
    slide.shapes.add_picture(img_y_0701, x0, y_top, width=img_width_y, height=img_height_y)
    slide.shapes.add_picture(img_y_0703, x0 + img_width_y + gap, y_top, width=img_width_y, height=img_height_y)

# 插入样本
for i in range(0, len(samples), 2):
    slide = prs.slides.add_slide(blank_slide)
    insert_sample(slide, samples[i], direction='left')
    if i + 1 < len(samples):
        insert_sample(slide, samples[i + 1], direction='right')

# 保存
prs.save(output_pptx)
print(f"✅ 已保存：{output_pptx}")
